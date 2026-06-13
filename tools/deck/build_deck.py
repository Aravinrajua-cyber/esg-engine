"""Standalone ESG Momentum Engine slide deck generator."""

from __future__ import annotations

import argparse
from pathlib import Path
from typing import Any

import yaml
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


TOOL_DIR = Path(__file__).resolve().parent
DEFAULT_CONTENT = TOOL_DIR / "deck_content.yaml"
DEFAULT_OUTPUT = TOOL_DIR / "ESG_Momentum_Engine.pptx"

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

BACKGROUND = RGBColor(0xFA, 0xFA, 0xF8)
TEXT = RGBColor(0x0B, 0x0B, 0x0C)
MUTED = RGBColor(0x66, 0x66, 0x63)
ACCENT = RGBColor(0x3B, 0x3B, 0xFF)

VALID_LAYOUTS = {"title", "statement", "hero_figure", "split", "closing"}
FIGURE_LAYOUTS = {"hero_figure", "split"}


class DeckContentError(ValueError):
    """Raised when the YAML content contract is invalid."""


def load_deck_content(path: Path) -> dict[str, Any]:
    with path.open(encoding="utf-8") as handle:
        raw = yaml.safe_load(handle)
    if not isinstance(raw, dict) or not isinstance(raw.get("slides"), list):
        raise DeckContentError("deck_content.yaml must contain a top-level slides list")
    for index, slide in enumerate(raw["slides"], start=1):
        if not isinstance(slide, dict):
            raise DeckContentError(f"slides[{index}] must be a mapping")
        if slide.get("layout") not in VALID_LAYOUTS:
            raise DeckContentError(f"slides[{index}].layout must be one of {sorted(VALID_LAYOUTS)}")
        if not slide.get("title"):
            raise DeckContentError(f"slides[{index}].title is required")
    return raw


def figure_path_for(content_path: Path, figure_path: str | None) -> Path | None:
    if not figure_path:
        return None
    path = Path(figure_path)
    return path if path.is_absolute() else (content_path.parent / path).resolve()


def missing_figures(content_path: Path) -> list[Path]:
    content = load_deck_content(content_path)
    missing: list[Path] = []
    for slide in content["slides"]:
        if slide["layout"] in FIGURE_LAYOUTS:
            figure = figure_path_for(content_path, slide.get("figure_path"))
            if figure is None or not figure.exists():
                missing.append(figure or Path("<missing figure_path>"))
    return missing


def set_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BACKGROUND


def add_textbox(
    slide,
    text: str,
    x,
    y,
    w,
    h,
    *,
    size: int,
    color: RGBColor = TEXT,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
    line_spacing: float = 0.94,
) -> None:
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = text or ""
    run.font.name = "Inter"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_body(slide, body: str | list[str] | None, x, y, w, h, *, size: int = 19) -> None:
    if body is None:
        return
    text = "\n".join(str(item) for item in body) if isinstance(body, list) else str(body)
    add_textbox(slide, text, x, y, w, h, size=size, color=MUTED, line_spacing=1.08)


def add_accent_rule(slide, x, y, w) -> None:
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.035))
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()


def add_kicker(slide, label: str, index: int) -> None:
    add_accent_rule(slide, Inches(0.58), Inches(6.98), Inches(0.26))
    add_textbox(
        slide,
        f"{index:02d} / {label.upper()}",
        Inches(0.92),
        Inches(6.86),
        Inches(3.8),
        Inches(0.28),
        size=8,
        color=MUTED,
        bold=True,
    )


def add_picture_cover(slide, image_path: Path, x, y, w, h) -> None:
    with Image.open(image_path) as image:
        img_w, img_h = image.size
    image_ratio = img_w / img_h
    box_ratio = w / h
    if image_ratio > box_ratio:
        pic_h = h
        pic_w = int(pic_h * image_ratio)
        pic_x = int(x - (pic_w - w) / 2)
        pic_y = y
    else:
        pic_w = w
        pic_h = int(pic_w / image_ratio)
        pic_x = x
        pic_y = int(y - (pic_h - h) / 2)
    slide.shapes.add_picture(str(image_path), pic_x, pic_y, width=pic_w, height=pic_h)


def build_title_slide(prs: Presentation, slide_data: dict[str, Any], index: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_accent_rule(slide, Inches(0.72), Inches(0.82), Inches(1.6))
    add_textbox(slide, slide_data["title"], Inches(0.72), Inches(1.48), Inches(8.9), Inches(1.8), size=52, bold=True)
    add_textbox(slide, slide_data.get("subtitle", ""), Inches(0.78), Inches(3.38), Inches(7.4), Inches(0.6), size=20, color=MUTED)
    add_body(slide, slide_data.get("body"), Inches(0.78), Inches(4.38), Inches(7.8), Inches(1.4), size=18)
    add_kicker(slide, "ESG Momentum Engine", index)


def build_statement_slide(prs: Presentation, slide_data: dict[str, Any], index: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_textbox(slide, slide_data["title"], Inches(0.78), Inches(1.2), Inches(11.3), Inches(3.3), size=44, bold=True, line_spacing=0.9)
    add_body(slide, slide_data.get("body") or slide_data.get("subtitle"), Inches(0.84), Inches(4.95), Inches(8.8), Inches(1.0), size=19)
    add_kicker(slide, "statement", index)


def build_hero_figure_slide(prs: Presentation, slide_data: dict[str, Any], index: int, content_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    figure = figure_path_for(content_path, slide_data.get("figure_path"))
    if figure is None or not figure.exists():
        raise FileNotFoundError(f"Missing figure for slide {index}: {slide_data.get('figure_path')}")
    add_picture_cover(slide, figure, 0, 0, SLIDE_W, SLIDE_H)
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(5.15), SLIDE_H)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = BACKGROUND
    overlay.line.fill.background()
    add_textbox(slide, slide_data["title"], Inches(0.58), Inches(0.72), Inches(4.0), Inches(1.5), size=28, bold=True)
    add_body(slide, slide_data.get("body") or slide_data.get("subtitle"), Inches(0.62), Inches(2.45), Inches(3.9), Inches(2.2), size=16)
    add_kicker(slide, "figure", index)


def build_split_slide(prs: Presentation, slide_data: dict[str, Any], index: int, content_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    figure = figure_path_for(content_path, slide_data.get("figure_path"))
    if figure is None or not figure.exists():
        raise FileNotFoundError(f"Missing figure for slide {index}: {slide_data.get('figure_path')}")
    add_textbox(slide, slide_data["title"], Inches(0.72), Inches(0.78), Inches(5.0), Inches(1.6), size=31, bold=True)
    add_body(slide, slide_data.get("body") or slide_data.get("subtitle"), Inches(0.76), Inches(2.65), Inches(4.8), Inches(2.7), size=18)
    add_picture_cover(slide, figure, Inches(6.2), Inches(0.54), Inches(6.55), Inches(5.92))
    add_kicker(slide, "split", index)


def build_closing_slide(prs: Presentation, slide_data: dict[str, Any], index: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_accent_rule(slide, Inches(0.72), Inches(1.08), Inches(1.6))
    add_textbox(slide, slide_data["title"], Inches(0.72), Inches(1.72), Inches(10.8), Inches(2.2), size=46, bold=True, line_spacing=0.92)
    add_body(slide, slide_data.get("body") or slide_data.get("subtitle"), Inches(0.78), Inches(4.4), Inches(8.6), Inches(1.2), size=20)
    add_kicker(slide, "closing", index)


def build_deck(content_path: Path = DEFAULT_CONTENT, output_path: Path = DEFAULT_OUTPUT) -> Path:
    content_path = content_path.resolve()
    output_path = output_path.resolve()
    content = load_deck_content(content_path)
    missing = missing_figures(content_path)
    if missing:
        raise FileNotFoundError("Missing figure(s): " + ", ".join(str(path) for path in missing))

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for index, slide_data in enumerate(content["slides"], start=1):
        layout = slide_data["layout"]
        if layout == "title":
            build_title_slide(prs, slide_data, index)
        elif layout == "statement":
            build_statement_slide(prs, slide_data, index)
        elif layout == "hero_figure":
            build_hero_figure_slide(prs, slide_data, index, content_path)
        elif layout == "split":
            build_split_slide(prs, slide_data, index, content_path)
        elif layout == "closing":
            build_closing_slide(prs, slide_data, index)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def main() -> int:
    parser = argparse.ArgumentParser(description="Generate ESG_Momentum_Engine.pptx from deck_content.yaml")
    parser.add_argument("--content", type=Path, default=DEFAULT_CONTENT, help="Path to deck_content.yaml")
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT, help="Output pptx path")
    args = parser.parse_args()
    print(build_deck(args.content, args.output))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
